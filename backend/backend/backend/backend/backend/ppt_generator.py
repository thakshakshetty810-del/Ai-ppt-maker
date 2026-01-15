from pptx import Presentation

def create_ppt(topic, purpose, slides, tone):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = f"{purpose} presentation ({tone})"

    # Content slides
    for i in range(slides - 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i + 1}"
        slide.placeholders[1].text = "AI-generated content will appear here"

    file_name = "generated_ppt.pptx"
    prs.save(file_name)

    return file_name
